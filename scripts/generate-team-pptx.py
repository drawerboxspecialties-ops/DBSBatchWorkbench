"""Generate a simple team PowerPoint for OptiCut CSV Splitter."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "public" / "DBS-OptiCut-Team-Overview.pptx"

# Brand-ish colors (teal / slate — matches station accent, not purple AI default)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_DARK = RGBColor(0x0F, 0x76, 0x6E)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF0, 0xFD, 0xFA)
LINE = RGBColor(0xCB, 0xD5, 0xE1)


def set_run(run, text, size=18, bold=False, color=INK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_bg(slide, color: RGBColor):
    fill = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    fill.line.fill.background()
    fill.fill.solid()
    fill.fill.fore_color.rgb = color
    spTree = slide.shapes._spTree
    sp = fill._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, TEAL_DARK)

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    set_run(p.add_run(), "DBS Drawers", size=40, bold=True, color=WHITE)
    p.alignment = PP_ALIGN.LEFT

    box2 = slide.shapes.add_textbox(Inches(0.9), Inches(3.3), Inches(11.5), Inches(1))
    tf2 = box2.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    set_run(p2.add_run(), "OptiCut CSV Splitter", size=28, bold=False, color=RGBColor(0x99, 0xF6, 0xE4))

    box3 = slide.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(11.5), Inches(1.2))
    tf3 = box3.text_frame
    tf3.word_wrap = True
    tf3.clear()
    p3 = tf3.paragraphs[0]
    set_run(
        p3.add_run(),
        "Allmoxy → batches → cut lists → OptiCut + Station\nSimple overview for the team",
        size=18,
        color=WHITE,
    )

    foot = slide.shapes.add_textbox(Inches(0.9), Inches(6.7), Inches(11), Inches(0.4))
    tf = foot.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    set_run(p.add_run(), "Drawer Box Specialties  ·  Internal tooling", size=12, color=RGBColor(0xCC, 0xFB, 0xF1))


def section_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    set_run(p.add_run(), text, size=28, bold=True, color=INK)


def bullets(slide, items, top=1.2, left=0.7, width=11.8, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        set_run(p.add_run(), "•  " + item, size=size, color=INK)


def content_slide(prs, title, items, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_accent_bar(slide)
    section_title(slide, title)
    bullets(slide, items)
    if note:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(11.8), Inches(0.5))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        set_run(p.add_run(), note, size=13, color=MUTED)
    return slide


def two_col_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_accent_bar(slide)
    section_title(slide, title)

    # left card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.2), Inches(5.8), Inches(5.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = SOFT
    card.line.color.rgb = RGBColor(0x99, 0xF6, 0xE4)

    lt = slide.shapes.add_textbox(Inches(0.85), Inches(1.4), Inches(5.3), Inches(0.5))
    tf = lt.text_frame
    tf.clear()
    set_run(tf.paragraphs[0].add_run(), left_title, size=18, bold=True, color=TEAL_DARK)
    bullets(slide, left_items, top=2.0, left=0.85, width=5.3, size=15)

    # right card
    card2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.2)
    )
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    card2.line.color.rgb = LINE

    rt = slide.shapes.add_textbox(Inches(7.05), Inches(1.4), Inches(5.3), Inches(0.5))
    tf = rt.text_frame
    tf.clear()
    set_run(tf.paragraphs[0].add_run(), right_title, size=18, bold=True, color=INK)
    bullets(slide, right_items, top=2.0, left=7.05, width=5.3, size=15)
    return slide


def flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_accent_bar(slide)
    section_title(slide, "Daily flow (simple)")

    steps = [
        ("1", "Import", "Drop Allmoxy\nCSV in the app"),
        ("2", "Review", "Check batches\n& specials"),
        ("3", "Print", "Cut lists +\nBatch index"),
        ("4", "Send", "Send batches\nto Station"),
        ("5", "Cut", "Floor scans &\nchecks lines"),
        ("6", "Export", "CSV / ZIP →\nOptiCut"),
    ]
    x0 = 0.55
    w = 1.9
    gap = 0.15
    for i, (num, label, detail) in enumerate(steps):
        x = x0 + i * (w + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(w), Inches(3.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = SOFT if i % 2 == 0 else RGBColor(0xF8, 0xFA, 0xFC)
        shape.line.color.rgb = RGBColor(0x99, 0xF6, 0xE4)

        nbox = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.2), Inches(w - 0.2), Inches(0.6))
        tf = nbox.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), num, size=28, bold=True, color=TEAL)

        lbox = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.9), Inches(w - 0.2), Inches(0.5))
        tf = lbox.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), label, size=16, bold=True, color=INK)

        dbox = slide.shapes.add_textbox(Inches(x + 0.1), Inches(3.5), Inches(w - 0.2), Inches(1.8))
        tf = dbox.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), detail, size=13, color=MUTED)

    note = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12), Inches(0.5))
    tf = note.text_frame
    tf.clear()
    set_run(
        tf.paragraphs[0].add_run(),
        "Prep computer does steps 1–4 & 6. Floor computer uses Station for step 5.",
        size=14,
        color=MUTED,
    )


def links_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_accent_bar(slide)
    section_title(slide, "Where to open it")

    rows = [
        ("App (prep)", "https://drawerboxspecialties-ops.github.io/OpticutExportAppV2/"),
        ("Station (floor)", "…/OpticutExportAppV2/#station   or tap Station in the header"),
        ("Guide", "Tap Guide in the app header (one-page reference)"),
    ]
    y = 1.5
    for label, url in rows:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(11.8), Inches(1.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = SOFT
        card.line.color.rgb = RGBColor(0x99, 0xF6, 0xE4)

        t = slide.shapes.add_textbox(Inches(0.95), Inches(y + 0.2), Inches(11.3), Inches(0.4))
        tf = t.text_frame
        tf.clear()
        set_run(tf.paragraphs[0].add_run(), label, size=16, bold=True, color=TEAL_DARK)

        u = slide.shapes.add_textbox(Inches(0.95), Inches(y + 0.6), Inches(11.3), Inches(0.4))
        tf = u.text_frame
        tf.clear()
        set_run(tf.paragraphs[0].add_run(), url, size=15, color=INK)
        y += 1.45


def closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, TEAL_DARK)

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(1))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), "Questions?", size=40, bold=True, color=WHITE)

    box2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10.3), Inches(1.5))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    set_run(
        p2.add_run(),
        "Goal: less spreadsheet work, clearer cut lists,\nfewer mistakes on the saw — every day.",
        size=18,
        color=RGBColor(0xCC, 0xFB, 0xF1),
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)

    content_slide(
        prs,
        "The problem",
        [
            "Allmoxy gives one big CSV for every order and part",
            "OptiCut needs smaller files by material and edge",
            "The floor needs clear cut lists with box counts",
            "Doing this by hand in spreadsheets is slow and easy to mess up",
        ],
        note="At ~800 boxes/day, small mistakes cost real time and material.",
    )

    content_slide(
        prs,
        "What this app does",
        [
            "Runs in the browser — drop in the Allmoxy CSV",
            "Splits work into batches (material + edge + ship date)",
            "Separates special orders (scoop, slope, dividers, etc.)",
            "Prints cut lists and a batch index with barcodes",
            "Exports OptiCut-ready CSV / ZIP",
            "Sends live cut lists to the Station screen on the floor",
        ],
    )

    flow_slide(prs)

    two_col_slide(
        prs,
        "What you get",
        "Prep / office",
        [
            "Split batches in the sidebar",
            "Cut-list print (landscape)",
            "Batch index + barcodes",
            "CSV or ZIP for OptiCut",
            "Send to Station",
        ],
        "Floor / Station",
        [
            "Live queue of sent batches",
            "Scan barcode to open a batch",
            "Check off lines as you cut",
            "Remove = hide (Add back later)",
            "Wipe database = password dbs",
        ],
    )

    content_slide(
        prs,
        "Box math (one rule)",
        [
            "Boxes = ceil(parts ÷ 4)  — same shop rule as always",
            "When GroupID exists: count boxes per group, then add them up",
            "Example: group 1 has 5 parts → 2 boxes; group 2 has 3 parts → 1 box → order = 3 boxes",
            "Print shows: 3 boxes (1-2, 2-1) next to the order",
            "Batch index shows the same next to each order number",
        ],
        note="Order totals and batch totals use this rule. Line Bx on the sheet is per cut line.",
    )

    content_slide(
        prs,
        "Batch names (quick read)",
        [
            "Format: CATEGORY_EDGE_firstOrder",
            "Examples: PLY_PVC_602480   ·   SPECIAL_PLY_CFB_602627",
            "PLY / FAA / SLD / MDF = material family",
            "PVC / CFB / RAW / TPE = top edge type",
            "SPECIAL_ = special ops (kept separate from normal runs)",
            "Full material + ship date show under the name on screen / print",
        ],
    )

    content_slide(
        prs,
        "Operator checklist",
        [
            "1. Export CSV from Allmoxy",
            "2. Open the app → drop the file → review batches",
            "3. Print Cut list(s) and Batch index",
            "4. Send to Station (or Send all)",
            "5. Export CSV / ZIP → import into OptiCut",
            "6. Floor: open Station → scan barcode → check lines",
        ],
    )

    links_slide(prs)
    closing_slide(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
